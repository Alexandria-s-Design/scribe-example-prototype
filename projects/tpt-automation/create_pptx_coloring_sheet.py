#!/usr/bin/env python3
"""
Create Ecosystem Coloring Sheet as PowerPoint (PPTX)
Then upload to Google Drive - PPTX has proper image sizing control
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Configuration
IMAGE_DIR = "C:/Users/MarieLexisDad/projects/tpt-automation/ecosystem-artwork"
OUTPUT_FILE = "C:/Users/MarieLexisDad/projects/tpt-automation/Ecosystem_Coloring_Sheet.pptx"

# Images with their positions (4 across top, 4 across bottom)
IMAGES = [
    {"file": "oak_tree.png", "left": 0.5, "top": 2.0, "name": "Oak Tree"},
    {"file": "deer.png", "left": 2.25, "top": 2.0, "name": "Deer"},
    {"file": "rabbit.png", "left": 4.0, "top": 2.0, "name": "Rabbit"},
    {"file": "mushrooms.png", "left": 5.75, "top": 2.0, "name": "Mushrooms"},
    {"file": "sun.png", "left": 0.5, "top": 4.5, "name": "Sun"},
    {"file": "rocks.png", "left": 2.25, "top": 4.5, "name": "Rocks"},
    {"file": "pond.png", "left": 4.0, "top": 4.5, "name": "Pond"},
    {"file": "grass_flowers.png", "left": 5.75, "top": 4.5, "name": "Grass"}
]

IMAGE_SIZE = 1.5  # inches

def create_coloring_sheet():
    """Create PowerPoint coloring sheet with proper sizing"""

    print("="*80)
    print("CREATING POWERPOINT COLORING SHEET")
    print("="*80)

    # Create presentation
    prs = Presentation()

    # Set slide size to US Letter (8.5" × 11")
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    print("\nAdding text elements...")

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(7.5), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Ecosystem Coloring Sheet - 3rd Grade"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Student info
    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(7.5), Inches(0.3)
    )
    info_frame = info_box.text_frame
    info_frame.text = "Name: ________________     Date: __________"
    info_para = info_frame.paragraphs[0]
    info_para.font.size = Pt(12)

    # Instructions
    instr_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6),
        Inches(7.5), Inches(0.3)
    )
    instr_frame = instr_box.text_frame
    instr_frame.text = "Instructions: Color each ecosystem element below"
    instr_para = instr_frame.paragraphs[0]
    instr_para.font.size = Pt(12)

    # Vocabulary footer
    vocab_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(7.5), Inches(0.4)
    )
    vocab_frame = vocab_box.text_frame
    vocab_frame.text = "Vocabulary: Ecosystem • Producer • Consumer • Decomposer"
    vocab_para = vocab_frame.paragraphs[0]
    vocab_para.font.size = Pt(11)
    vocab_para.font.italic = True

    print("Adding ecosystem images...")

    # Add images
    for img_info in IMAGES:
        img_path = os.path.join(IMAGE_DIR, img_info['file'])

        if not os.path.exists(img_path):
            print(f"  WARNING: {img_info['file']} not found")
            continue

        print(f"  Adding: {img_info['name']}")

        # Add image with exact size and position
        slide.shapes.add_picture(
            img_path,
            Inches(img_info['left']),
            Inches(img_info['top']),
            width=Inches(IMAGE_SIZE),
            height=Inches(IMAGE_SIZE)
        )

        # Add label below image
        label_box = slide.shapes.add_textbox(
            Inches(img_info['left']),
            Inches(img_info['top'] + IMAGE_SIZE + 0.05),
            Inches(IMAGE_SIZE),
            Inches(0.25)
        )
        label_frame = label_box.text_frame
        label_frame.text = img_info['name']
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(9)
        label_para.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(OUTPUT_FILE)

    print(f"\n PowerPoint file created: {OUTPUT_FILE}")
    print(f"File size: {os.path.getsize(OUTPUT_FILE) / 1024:.1f} KB")

    print("\n" + "="*80)
    print("CREATION COMPLETE")
    print("="*80)
    print("\nYou can now:")
    print("1. Open the PPTX file directly in PowerPoint")
    print("2. Upload to Google Drive and convert to Google Slides")
    print("3. Export as PDF for printing")

    return OUTPUT_FILE

if __name__ == "__main__":
    try:
        output = create_coloring_sheet()
        print(f"\nOutput file: {output}")
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()
