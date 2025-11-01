#!/usr/bin/env python3
"""
Create Beautiful Word Search in PowerPoint
Ecosystem-themed for 3rd graders with professional design
"""

import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration
OUTPUT_FILE = "C:/Users/MarieLexisDad/projects/tpt-automation/Ecosystem_Word_Search.pptx"

# 10 ecosystem words for 3rd grade
WORDS = [
    "ECOSYSTEM",
    "PRODUCER",
    "CONSUMER",
    "DECOMPOSER",
    "HABITAT",
    "ENERGY",
    "FOOD",
    "CHAIN",
    "PLANT",
    "ANIMAL"
]

GRID_SIZE = 15  # 15x15 grid

class WordSearch:
    def __init__(self, words, grid_size):
        self.words = words
        self.grid_size = grid_size
        self.grid = [['' for _ in range(grid_size)] for _ in range(grid_size)]
        self.placed_words = []

    def can_place_word(self, word, row, col, direction):
        """Check if word can be placed at position in given direction"""
        dr, dc = direction

        for i, letter in enumerate(word):
            r, c = row + i * dr, col + i * dc

            # Check bounds
            if r < 0 or r >= self.grid_size or c < 0 or c >= self.grid_size:
                return False

            # Check if cell is empty or matches
            if self.grid[r][c] != '' and self.grid[r][c] != letter:
                return False

        return True

    def place_word(self, word, row, col, direction):
        """Place word at position in given direction"""
        dr, dc = direction
        positions = []

        for i, letter in enumerate(word):
            r, c = row + i * dr, col + i * dc
            self.grid[r][c] = letter
            positions.append((r, c))

        self.placed_words.append({
            'word': word,
            'positions': positions,
            'direction': direction
        })

    def generate(self):
        """Generate the word search puzzle"""
        # Directions: right, down, diagonal down-right
        directions = [(0, 1), (1, 0), (1, 1)]

        random.shuffle(self.words)

        for word in self.words:
            placed = False
            attempts = 0
            max_attempts = 100

            while not placed and attempts < max_attempts:
                # Random position
                row = random.randint(0, self.grid_size - 1)
                col = random.randint(0, self.grid_size - 1)
                direction = random.choice(directions)

                if self.can_place_word(word, row, col, direction):
                    self.place_word(word, row, col, direction)
                    placed = True

                attempts += 1

            if not placed:
                print(f"Warning: Could not place word '{word}'")

        # Fill empty cells with random letters
        for r in range(self.grid_size):
            for c in range(self.grid_size):
                if self.grid[r][c] == '':
                    self.grid[r][c] = random.choice('ABCDEFGHIJKLMNOPQRSTUVWXYZ')

        return self.grid

def create_word_search_pptx():
    """Create beautiful word search PowerPoint"""

    print("="*80)
    print("CREATING ECOSYSTEM WORD SEARCH")
    print("="*80)

    # Generate word search
    print("\nGenerating puzzle...")
    ws = WordSearch(WORDS, GRID_SIZE)
    grid = ws.generate()

    print(f"Placed {len(ws.placed_words)}/{len(WORDS)} words")

    # Create PowerPoint
    print("\nCreating PowerPoint...")
    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add decorative header background
    header_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(8.5), Inches(1.2)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(34, 139, 34)  # Forest green
    header_bg.line.color.rgb = RGBColor(34, 139, 34)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(7.5), Inches(0.4)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🌳 Ecosystem Word Search 🌳"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7),
        Inches(7.5), Inches(0.3)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Find all 10 words hidden in the grid!"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(14)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Student info
    info_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(7.5), Inches(0.3)
    )
    info_frame = info_box.text_frame
    info_frame.text = "Name: ________________     Date: __________     Grade: 3rd"
    info_para = info_frame.paragraphs[0]
    info_para.font.size = Pt(12)
    info_para.alignment = PP_ALIGN.CENTER

    # Word search grid
    # Calculate cell size for perfect square
    grid_width = 6.0  # inches
    cell_size = grid_width / GRID_SIZE
    grid_left = (8.5 - grid_width) / 2  # Center horizontally
    grid_top = 2.0

    print(f"Creating {GRID_SIZE}×{GRID_SIZE} grid...")

    # Add grid cells
    for row in range(GRID_SIZE):
        for col in range(GRID_SIZE):
            # Cell background
            cell = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(grid_left + col * cell_size),
                Inches(grid_top + row * cell_size),
                Inches(cell_size),
                Inches(cell_size)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Honeydew
            cell.line.color.rgb = RGBColor(34, 139, 34)  # Forest green
            cell.line.width = Pt(0.5)

            # Letter
            letter_box = slide.shapes.add_textbox(
                Inches(grid_left + col * cell_size),
                Inches(grid_top + row * cell_size),
                Inches(cell_size),
                Inches(cell_size)
            )
            letter_frame = letter_box.text_frame
            letter_frame.text = grid[row][col]
            letter_frame.vertical_anchor = 1  # Middle
            letter_para = letter_frame.paragraphs[0]
            letter_para.font.size = Pt(11)
            letter_para.font.bold = True
            letter_para.font.color.rgb = RGBColor(0, 0, 0)
            letter_para.alignment = PP_ALIGN.CENTER

    # Word list section
    word_list_top = grid_top + GRID_SIZE * cell_size + 0.3

    # Word list header background
    word_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(word_list_top),
        Inches(7.5), Inches(0.4)
    )
    word_bg.fill.solid()
    word_bg.fill.fore_color.rgb = RGBColor(34, 139, 34)
    word_bg.line.color.rgb = RGBColor(34, 139, 34)

    # Word list header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(word_list_top + 0.05),
        Inches(7.5), Inches(0.3)
    )
    header_frame = header_box.text_frame
    header_frame.text = "Find These Words:"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(14)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER

    # Word list (2 columns)
    words_sorted = sorted(WORDS)
    words_per_col = 5
    col_width = 3.5

    for i, word in enumerate(words_sorted):
        col_idx = i // words_per_col
        row_idx = i % words_per_col

        # Checkbox + word
        word_left = 1.0 + col_idx * col_width
        word_top = word_list_top + 0.5 + row_idx * 0.35

        # Checkbox
        checkbox = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(word_left), Inches(word_top),
            Inches(0.2), Inches(0.2)
        )
        checkbox.fill.solid()
        checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        checkbox.line.color.rgb = RGBColor(0, 0, 0)
        checkbox.line.width = Pt(1)

        # Word text
        word_box = slide.shapes.add_textbox(
            Inches(word_left + 0.3), Inches(word_top - 0.03),
            Inches(2.5), Inches(0.25)
        )
        word_frame = word_box.text_frame
        word_frame.text = word
        word_para = word_frame.paragraphs[0]
        word_para.font.size = Pt(12)
        word_para.font.bold = False

    # Footer with instructions
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(10.3),
        Inches(7.5), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Tip: Words can go across (→), down (↓), or diagonal (↘)"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(10)
    footer_para.font.italic = True
    footer_para.alignment = PP_ALIGN.CENTER

    # Save
    prs.save(OUTPUT_FILE)

    print(f"\nPowerPoint created: {OUTPUT_FILE}")
    print(f"File size: {os.path.getsize(OUTPUT_FILE) / 1024:.1f} KB")

    print("\n" + "="*80)
    print("WORD SEARCH COMPLETE")
    print("="*80)
    print("\nWords included:")
    for word in sorted(WORDS):
        print(f"  ✓ {word}")

    return OUTPUT_FILE

if __name__ == "__main__":
    try:
        create_word_search_pptx()
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()
